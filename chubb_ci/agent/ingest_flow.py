"""C1 — Document ingestion: competitor decks/brochures → database, with provenance.

Two lanes, per the Verify design:
- **Structured lane (deterministic)**: per-brand product tables are parsed by
  `tools/pptx_ingest` (no LLM) and applied directly — structured parses carry high
  base confidence.
- **Unstructured lane (LLM)**: brand-profile slides (品牌基本情况/产品分析) are
  LLM-extracted into per-field claims. Every claim carries its slide reference as
  source and goes through 真实性核查; single-source LLM claims fall below the
  auto-accept threshold by design, so they queue for human 采纳/驳回 on the dashboard
  rather than silently overwriting curated brand profiles.
"""

from __future__ import annotations

import json
import re
from pathlib import Path

from chubb_ci.agent.runtime import AgentContext, finish_run
from chubb_ci.agent.state import SourcedFact
from chubb_ci.agent.verify import BASE_LLM, verify_facts
from chubb_ci.llm.base import LLMClient, LLMError
from chubb_ci.llm.factory import resolve_model
from chubb_ci.tools.pptx_ingest import BRAND_MAP, ingest_pptx

_PROFILE_FIELDS = {
    "positioning": "品牌定位",
    "competition_tier": "竞争层级",
    "target_audience": "核心客群",
    "market_scale": "规模与销量",
    "supply_chain": "供应链与产地",
    "warranty": "售后质保",
}

_PROFILE_SYSTEM = """你是竞争情报分析助手。从竞品分析幻灯片的文字中，提取该品牌的档案信息。
只输出一个 JSON 对象，可包含键：positioning, competition_tier, target_audience,
market_scale, supply_chain, warranty（值为中文短句）。**页面中没有的信息不要输出对应键，
严禁推断或编造。** 不要输出解释或 Markdown。"""

_TITLE_PROFILE_RE = re.compile(r"^(.+?)(?:品牌)?(基本情况|产品分析)$")


def _profile_slides(path: Path) -> dict[str, list[tuple[int, str]]]:
    """Collect {brand: [(slide_no, slide_text)]} for 基本情况/产品分析 slides."""
    from pptx import Presentation

    prs = Presentation(str(path))
    out: dict[str, list[tuple[int, str]]] = {}
    for idx, slide in enumerate(prs.slides, 1):
        texts = [s.text_frame.text.strip() for s in slide.shapes
                 if s.has_text_frame and s.text_frame.text.strip()]
        if not texts:
            continue
        title = texts[0].splitlines()[0].strip()
        m = _TITLE_PROFILE_RE.match(title)
        if not m:
            continue
        raw_brand = m.group(1).strip()
        brand = BRAND_MAP.get(raw_brand)
        if brand is None:
            # try suffix matching ("美国迪堡" → "迪堡")
            for key, canon in BRAND_MAP.items():
                if key in raw_brand:
                    brand = canon
                    break
        if brand is None:
            continue
        body = "\n".join(texts)[:6000]
        out.setdefault(brand, []).append((idx, body))
    return out


def _extract_profile_facts(
    ctx: AgentContext, llm: LLMClient, brand: str,
    slides: list[tuple[int, str]], doc_name: str,
) -> list[SourcedFact]:
    facts: list[SourcedFact] = []
    for slide_no, text in slides:
        source = f"{doc_name}#slide{slide_no}"
        try:
            resp = ctx.track(llm.complete(
                system=_PROFILE_SYSTEM,
                user=f"品牌：{brand}\n\n幻灯片文字：\n<<<\n{text}\n>>>",
                model=resolve_model(ctx.settings, "extract"),
                json_mode=True,
            ))
            start, end = resp.content.find("{"), resp.content.rfind("}")
            data = json.loads(resp.content[start:end + 1]) if start != -1 else {}
        except (LLMError, json.JSONDecodeError, ValueError) as exc:
            ctx.log("抽取", f"{brand} slide{slide_no} 档案抽取失败：{exc}")
            continue
        for field, label in _PROFILE_FIELDS.items():
            value = data.get(field)
            if not (value and isinstance(value, str) and value.strip()):
                continue
            # Drop non-informative extractions the prompt forbids but LLMs still emit.
            if value.strip() in ("未提及", "无", "不详", "未知", "N/A", "-", "暂无"):
                continue
            facts.append(SourcedFact(
                    claim=f"{brand} 的{label}：{value.strip()}",
                    field=field, value=value.strip(), subject=brand,
                    sources=[source], confidence=BASE_LLM,
                ))
    return facts


def run_ingest(ctx: AgentContext, llm: LLMClient, path: str | Path) -> None:
    path = Path(path)
    if not path.exists():
        finish_run(ctx, status="failed", error=f"file not found: {path}")
        return
    doc_name = path.name
    ctx.log("规划", f"开始摄取文档：{doc_name}", f"{path.stat().st_size/1e6:.1f} MB")

    # --- Structured lane: deterministic product tables --------------------
    if path.suffix.lower() == ".pptx":
        result = ingest_pptx(ctx.session, path)
        ctx.log("抽取", f"结构化表格：{result['brands']} 个品牌 {result['products']} 款产品（确定性解析，已入库）")

        # --- Unstructured lane: brand-profile slides via LLM ---------------
        profiles = _profile_slides(path)
        ctx.log("抽取", f"发现 {len(profiles)} 个品牌的档案/分析页，开始 LLM 抽取（每条信息附带来源页）")
        facts: list[SourcedFact] = []
        for brand, slides in profiles.items():
            facts.extend(_extract_profile_facts(ctx, llm, brand, slides, doc_name))
            ctx.check_budget()
        ctx.log("抽取", f"共提取 {len(facts)} 条品牌档案声明")
    else:
        # PDF / other: whole-text best effort via the pdf extractor.
        from chubb_ci.tools.pdf_ingest import extract_pdf_text

        text = extract_pdf_text(path)[:12000]
        ctx.log("抽取", f"PDF 文本 {len(text)} 字，LLM 抽取品牌档案声明")
        facts = []
        for brand in set(BRAND_MAP.values()):
            if brand.split()[0] in text:
                facts.extend(_extract_profile_facts(
                    ctx, llm, brand, [(0, text)], doc_name))
        result = {"brands": 0, "products": 0}

    # --- Verify (真实性核查) ------------------------------------------------
    vr = verify_facts(facts, threshold=ctx.settings.agent_verify_threshold)
    ctx.run.facts_total = len(facts)
    ctx.run.facts_verified = len(vr.verified)
    ctx.run.facts_pending = len(vr.pending)
    ctx.log("核查", f"核查完成：{len(vr.verified)} 条自动通过，{len(vr.pending)} 条待人工确认",
            "单一来源的 LLM 抽取默认进入人工确认队列")
    for f in vr.pending + vr.verified:
        ctx.add_pending_fact(
            subject=f.subject, field=f.field, value=f.value, claim=f.claim,
            sources=f.sources, corroborations=f.corroborations,
            confidence=f.confidence, status=f.status,
        )

    # --- Report -------------------------------------------------------------
    lines = [f"# 文档摄取报告：{doc_name}", "",
             f"- 结构化产品表：**{result['products']}** 款（{result['brands']} 个品牌，已直接入库）",
             f"- 品牌档案声明：**{len(facts)}** 条（自动通过 {len(vr.verified)} / 待人工 {len(vr.pending)}）",
             "", "## 待人工确认的声明（前 10 条）", ""]
    for f in vr.pending[:10]:
        lines.append(f"- {f.claim}（来源：{'、'.join(f.sources)}，置信 {f.confidence}）")
    ctx.log("应用", "摄取完成，待确认声明已进入审核队列")
    finish_run(ctx, result_md="\n".join(lines))
