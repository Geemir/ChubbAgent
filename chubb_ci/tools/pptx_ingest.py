"""Ingest competitor product data from the marketing deck (保险柜市场竞品分析 PPTX).

The deck carries per-brand product tables (slides titled "XX产品（热卖/功能贴近/价值最高）")
in key-value form: 系列名 / 产品情况 / 销量 / 价格 / 认证 / 规格 / 备注, with one column per
product. Parsing is fully deterministic (no LLM): prices like ``￥7,500~26,000+`` take the
range low end; sales like ``TB:2W+ JD:1W+`` sum with 万(W)=10000; the 认证 text is routed
to fire/euro/GB fields and standardized via chubb_ci/normalize.

Records land as ProductRecords under channel ``分析报告`` (one synthetic snapshot per brand),
so the deck's real market data flows into products / market-map / benchmark. Re-running is
idempotent (previous 分析报告 snapshots for the deck are replaced).
"""

from __future__ import annotations

import re
from datetime import datetime, timezone
from pathlib import Path

from loguru import logger
from sqlmodel import Session, delete, select

from chubb_ci.diff.matching import normalize_product_key
from chubb_ci.normalize import fire_hours, security_score
from chubb_ci.schemas.models import ProductRecord, Snapshot, SnapshotStatus

# Slide-title brand → canonical brand name (aligned with config/brands.yaml).
BRAND_MAP: dict[str, str] = {
    "德国百卫特": "德国百卫特 Burg Wächter",
    "百卫特": "德国百卫特 Burg Wächter",
    "美国迪堡": "迪堡 Diebold",
    "迪堡": "迪堡 Diebold",
    "美国善卫": "善卫 SentrySafe",
    "善卫": "善卫 SentrySafe",
    "福美德": "福美德 Format",
    "艾斐堡": "艾斐堡 Afeibao",
    "艾谱": "艾谱 AIPU",
    "永发": "永发 Yongfa",
    "盾牌": "盾牌 Dunpai",
}

_TITLE_RE = re.compile(r"^(.+?)产品（")
_NUM_RE = re.compile(r"(\d[\d,]*(?:\.\d+)?)\s*([wW万])?")

_FIRE_KEYWORDS = ("防火", "耐火", "fire")
_EURO_KEYWORDS = ("S2", "EN", "欧标", "UL687", "级")
_GB_KEYWORDS = ("国标", "3C", "GB")


def _parse_price(text: str | None) -> float | None:
    """First numeric amount (range low end); handles ￥/RMB/commas/multi-variant text."""
    if not text:
        return None
    m = _NUM_RE.search(text.replace("，", ","))
    if not m:
        return None
    value = float(m.group(1).replace(",", ""))
    if m.group(2):  # 万
        value *= 10_000
    return value if value > 0 else None


def _parse_sales(text: str | None) -> int | None:
    """Sum all numbers in the cell; W/万 multiplies by 10000 ("TB:2W+ JD:1W+" → 30000)."""
    if not text:
        return None
    total = 0.0
    for m in _NUM_RE.finditer(text.replace("，", ",")):
        v = float(m.group(1).replace(",", ""))
        if m.group(2):
            v *= 10_000
        total += v
    return int(total) if total > 0 else None


def _route_cert(text: str | None) -> tuple[str | None, str | None, str | None]:
    """Split the 认证 cell into (fire_rating, euro_grade, gb_grade) raw texts."""
    if not text:
        return None, None, None
    s = text.strip()
    if not s or s in ("无", "-", "无防火防盗认证"):
        return None, None, None
    fire = s if any(k in s for k in _FIRE_KEYWORDS) or re.search(r"\d+\s*(分钟|min|h)", s, re.I) else None
    euro = s if any(k in s for k in _EURO_KEYWORDS) else None
    gb = s if any(k in s for k in _GB_KEYWORDS) else None
    return fire, euro, gb


def parse_deck(path: str | Path) -> dict[str, list[dict]]:
    """Parse the PPTX and return {brand: [product dicts]}."""
    from pptx import Presentation

    path = Path(path)
    if not path.exists():
        raise FileNotFoundError(path)
    prs = Presentation(str(path))

    out: dict[str, list[dict]] = {}
    for slide in prs.slides:
        title = ""
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                title = shape.text_frame.text.strip().splitlines()[0]
                break
        m = _TITLE_RE.match(title)
        if not m:
            continue
        brand = BRAND_MAP.get(m.group(1).strip(), m.group(1).strip())

        for shape in slide.shapes:
            if not shape.has_table:
                continue
            table = shape.table
            rows = [[c.text.strip() for c in r.cells] for r in table.rows]
            if not rows or not rows[0]:
                continue
            n_products = len(rows[0]) - 1
            for col in range(1, n_products + 1):
                fields = {r[0]: (r[col] if col < len(r) else "") for r in rows if r and r[0]}
                name = fields.get("系列名", "").strip()
                if not name:
                    continue
                fire, euro, gb = _route_cert(fields.get("认证"))
                features = [v for v in (fields.get("规格"), fields.get("备注")) if v and v != "无"]
                out.setdefault(brand, []).append({
                    "product_name": name,
                    "status_label": (fields.get("产品情况") or "").strip() or None,
                    "sales_volume": _parse_sales(fields.get("销量")),
                    "price": _parse_price(fields.get("价格")),
                    "fire_rating": fire,
                    "euro_grade": euro,
                    "gb_grade": gb,
                    "key_features": features,
                })
    return out


def ingest_pptx(session: Session, path: str | Path) -> dict:
    """Import the deck's product tables as 分析报告-channel ProductRecords."""
    data = parse_deck(path)
    now = datetime.now(timezone.utc)
    deck_name = Path(path).name

    total = 0
    for brand, products in data.items():
        source_name = f"pptx-{normalize_product_key(brand)}"

        # Idempotency: drop previous records/snapshots from this synthetic source.
        old_ids = [s.id for s in session.exec(
            select(Snapshot).where(Snapshot.source_name == source_name)).all()]
        if old_ids:
            session.exec(delete(ProductRecord).where(ProductRecord.snapshot_id.in_(old_ids)))  # type: ignore[union-attr]
            session.exec(delete(Snapshot).where(Snapshot.source_name == source_name))

        snap = Snapshot(
            source_name=source_name, company=brand, url=deck_name,
            page_type="pricing", channel="分析报告",
            status=SnapshotStatus.OK.value, crawl_time=now, num_products=len(products),
        )
        session.add(snap)
        session.commit()
        session.refresh(snap)

        for p in products:
            session.add(ProductRecord(
                snapshot_id=snap.id,
                source_name=source_name,
                company=brand,
                product_name=p["product_name"],
                product_key=normalize_product_key(p["product_name"]),
                category="保险柜",
                price=p["price"],
                currency="CNY",
                sales_volume=p["sales_volume"],
                status_label=p["status_label"],
                fire_rating=p["fire_rating"],
                euro_grade=p["euro_grade"],
                gb_grade=p["gb_grade"],
                key_features=p["key_features"],
                source_url=deck_name,
                crawl_time=now,
                fire_hours=fire_hours(p["fire_rating"]),
                security_score=security_score(p["gb_grade"], p["euro_grade"]),
            ))
            total += 1
        session.commit()

    logger.info("ingested {} products for {} brands from {}", total, len(data), deck_name)
    return {"brands": len(data), "products": total}
